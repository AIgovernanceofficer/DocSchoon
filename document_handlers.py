"""
Verwerking van bestanden: tekst eruit halen, anonimiseren en weer terugzetten
met behoud van opmaak.

Ondersteunde formaten:
  - .docx  (Word)        via python-docx
  - .pptx  (PowerPoint)  via python-pptx
  - .xlsx  (Excel)       via openpyxl
  - .txt   (platte tekst)

Kernidee voor Word/PowerPoint
-----------------------------
In .docx en .pptx is een alinea opgebouwd uit "runs": stukjes tekst met elk
hun eigen opmaak (lettertype, vet, kleur). Eén woord kan over meerdere runs
verdeeld zijn. Om gegevens betrouwbaar te vinden EN de opmaak te behouden:

  1. plak de tekst van alle runs in een alinea aan elkaar,
  2. detecteer de persoonsgegevens op die volledige alineatekst,
  3. zet de vervangingen terug in de juiste runs op basis van tekenposities.

Zo worden ook gegevens gevonden die over meerdere runs lopen, terwijl de
opmaak van de niet-geraakte runs intact blijft.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field

from .anonymizer import PIIAnonymizer


@dataclass
class ProcessResult:
    """Samenvatting van een verwerkt bestand."""

    output_path: str
    total_replacements: int = 0
    counts_by_type: dict[str, int] = field(default_factory=dict)

    def add(self, entity_type: str, n: int = 1) -> None:
        self.total_replacements += n
        self.counts_by_type[entity_type] = self.counts_by_type.get(entity_type, 0) + n


SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".txt"}


# ---------------------------------------------------------------------------
# Run-bewerking (gedeeld door Word en PowerPoint)
# ---------------------------------------------------------------------------
def _edit_runs(runs, anonymizer: PIIAnonymizer, style: str, result: ProcessResult) -> None:
    """
    Anonimiseer de tekst van een lijst runs en schrijf het resultaat terug.

    'runs' is een lijst objecten met een .text-attribuut (python-docx of
    python-pptx runs). De functie muteert run.text rechtstreeks.
    """
    full_text = "".join(run.text or "" for run in runs)
    if not full_text.strip():
        return

    detections = anonymizer.detect(full_text)
    if not detections:
        return

    # vervangingen + telling
    replacements: list[tuple[int, int, str]] = []
    for d in detections:
        from .config import label_for

        if style == "redact":
            repl = "█" * (d.end - d.start)
        else:
            repl = f"[{label_for(d.entity_type, anonymizer.language)}]"
        replacements.append((d.start, d.end, repl))
        result.add(d.entity_type)

    # bereken de tekenposities (offsets) van elke run in de samengevoegde tekst
    offsets: list[tuple[int, int]] = []
    pos = 0
    for run in runs:
        length = len(run.text or "")
        offsets.append((pos, pos + length))
        pos += length

    # verzamel per run de lokale bewerkingen
    run_edits: list[list[tuple[int, int, str]]] = [[] for _ in runs]
    for gs, ge, repl in replacements:
        first = True
        for i, (rs, re_) in enumerate(offsets):
            if re_ <= gs or rs >= ge:
                continue  # geen overlap met deze run
            local_start = max(gs, rs) - rs
            local_end = min(ge, re_) - rs
            insert = repl if first else ""  # vervangtekst in de eerste geraakte run
            run_edits[i].append((local_start, local_end, insert))
            first = False

    # pas de bewerkingen per run toe (rechts-naar-links binnen de run)
    for i, run in enumerate(runs):
        edits = sorted(run_edits[i], key=lambda x: x[0], reverse=True)
        if not edits:
            continue
        text = run.text or ""
        for local_start, local_end, insert in edits:
            text = text[:local_start] + insert + text[local_end:]
        run.text = text


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------
def _process_docx(in_path, out_path, anonymizer, style, result) -> None:
    from docx import Document

    doc = Document(in_path)

    def handle_paragraphs(paragraphs):
        for paragraph in paragraphs:
            _edit_runs(paragraph.runs, anonymizer, style, result)

    def handle_tables(tables):
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                    handle_paragraphs(cell.paragraphs)
                    handle_tables(cell.tables)  # geneste tabellen

    # hoofdtekst
    handle_paragraphs(doc.paragraphs)
    handle_tables(doc.tables)

    # kop- en voetteksten (namen staan vaak in briefhoofden/voetregels)
    for section in doc.sections:
        for hf in (section.header, section.footer,
                   section.first_page_header, section.first_page_footer,
                   section.even_page_header, section.even_page_footer):
            try:
                handle_paragraphs(hf.paragraphs)
                handle_tables(hf.tables)
            except Exception:
                pass

    doc.save(out_path)


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------
def _process_pptx(in_path, out_path, anonymizer, style, result) -> None:
    from pptx import Presentation

    prs = Presentation(in_path)

    def handle_text_frame(text_frame):
        for paragraph in text_frame.paragraphs:
            _edit_runs(paragraph.runs, anonymizer, style, result)

    def handle_shape(shape):
        if shape.has_text_frame:
            handle_text_frame(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    handle_text_frame(cell.text_frame)
        # gegroepeerde vormen: doorzoek de subvormen
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for sub in shape.shapes:
                handle_shape(sub)

    for slide in prs.slides:
        for shape in slide.shapes:
            handle_shape(shape)
        # notities bij de slide
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf is not None:
                handle_text_frame(notes_tf)

    prs.save(out_path)


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------
def _process_xlsx(in_path, out_path, anonymizer, style, result) -> None:
    from openpyxl import load_workbook

    wb = load_workbook(in_path)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value.strip():
                    new_value, detections = anonymizer.anonymize_text(cell.value, style)
                    if detections:
                        cell.value = new_value
                        for d in detections:
                            result.add(d.entity_type)
    wb.save(out_path)


# ---------------------------------------------------------------------------
# Platte tekst (.txt)
# ---------------------------------------------------------------------------
def _process_txt(in_path, out_path, anonymizer, style, result) -> None:
    with open(in_path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    new_text, detections = anonymizer.anonymize_text(text, style)
    for d in detections:
        result.add(d.entity_type)
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(new_text)


# ---------------------------------------------------------------------------
# Publieke ingang
# ---------------------------------------------------------------------------
_HANDLERS = {
    ".docx": _process_docx,
    ".pptx": _process_pptx,
    ".xlsx": _process_xlsx,
    ".txt": _process_txt,
}


def anonymize_file(
    input_path: str,
    output_path: str,
    anonymizer: PIIAnonymizer,
    style: str = "tag",
) -> ProcessResult:
    """
    Anonimiseer één bestand en schrijf het schone bestand naar output_path.

    Het bestandstype wordt afgeleid uit de extensie van input_path.
    """
    ext = os.path.splitext(input_path)[1].lower()
    handler = _HANDLERS.get(ext)
    if handler is None:
        raise ValueError(
            f"Niet-ondersteund bestandstype '{ext}'. "
            f"Ondersteund: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    result = ProcessResult(output_path=output_path)
    handler(input_path, output_path, anonymizer, style, result)
    return result
